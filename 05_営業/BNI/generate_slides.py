from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ===== カラー定義 =====
COLOR_BG       = RGBColor(0xFF, 0xFF, 0xFF)  # 白
COLOR_ACCENT   = RGBColor(0x1A, 0x73, 0xE8)  # ブルー
COLOR_DARK     = RGBColor(0x1A, 0x1A, 0x2E)  # ダークネイビー
COLOR_GRAY     = RGBColor(0x55, 0x55, 0x55)  # グレー
COLOR_LIGHT    = RGBColor(0xF0, 0xF4, 0xFF)  # 薄ブルー背景

W = Inches(13.33)  # 16:9 横幅
H = Inches(7.5)    # 16:9 縦幅


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # 完全ブランク
    return prs.slides.add_slide(layout)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=24, bold=False, color=COLOR_DARK,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # MSO_SHAPE_TYPE.RECTANGLE = 1
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def slide_title_only(slide, title, subtitle=None):
    """タイトル中央寄せレイアウト"""
    set_bg(slide, COLOR_BG)
    # アクセントバー（左）
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, COLOR_ACCENT)

    add_textbox(slide, title,
                Inches(1), Inches(2.8), Inches(11), Inches(1.5),
                font_size=44, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(1), Inches(4.4), Inches(11), Inches(0.8),
                    font_size=22, color=COLOR_GRAY, align=PP_ALIGN.CENTER)


def slide_section(slide, title, body_lines, note=None):
    """タイトル＋箇条書きレイアウト"""
    set_bg(slide, COLOR_BG)
    # ヘッダーバー
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.2), COLOR_DARK)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.2), Inches(12), Inches(0.8),
                font_size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    # 本文
    body_text = "\n".join(body_lines)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in body_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(22)
        run.font.color.rgb = COLOR_DARK
    if note:
        slide.notes_slide.notes_text_frame.text = note


def slide_comparison(slide, title, lead, rows, note=None):
    """左右比較表レイアウト"""
    set_bg(slide, COLOR_BG)
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.2), COLOR_DARK)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.2), Inches(12), Inches(0.8),
                font_size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if lead:
        add_textbox(slide, lead,
                    Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
                    font_size=20, italic=True, color=COLOR_GRAY)
    # 比較ヘッダー
    add_rect(slide, Inches(0.8), Inches(1.9), Inches(4.5), Inches(0.5), COLOR_LIGHT)
    add_rect(slide, Inches(7.0), Inches(1.9), Inches(4.5), Inches(0.5), RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide, "実店舗 2店舗目",
                Inches(0.8), Inches(1.9), Inches(4.5), Inches(0.5),
                font_size=18, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, "ネットショップ",
                Inches(7.0), Inches(1.9), Inches(4.5), Inches(0.5),
                font_size=18, bold=True, color=RGBColor(0x2E, 0x7D, 0x32), align=PP_ALIGN.CENTER)

    for i, (left_text, right_text) in enumerate(rows):
        y = Inches(2.55) + Inches(0.75) * i
        add_textbox(slide, left_text,  Inches(0.8), y, Inches(4.5), Inches(0.6),
                    font_size=20, color=COLOR_GRAY)
        add_textbox(slide, "→",  Inches(5.5), y, Inches(1.2), Inches(0.6),
                    font_size=20, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, right_text, Inches(7.0), y, Inches(4.5), Inches(0.6),
                    font_size=20, bold=True, color=RGBColor(0x2E, 0x7D, 0x32))
    if note:
        slide.notes_slide.notes_text_frame.text = note


def slide_contact(slide, title, customer_title, customer, bullets,
                  connect_title, connects, note=None):
    """コンタクトサークルレイアウト"""
    set_bg(slide, COLOR_BG)
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.2), COLOR_DARK)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.2), Inches(12), Inches(0.8),
                font_size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    # 上段：お客様
    add_rect(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.45), COLOR_LIGHT)
    add_textbox(slide, customer_title,
                Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.45),
                font_size=18, bold=True, color=COLOR_ACCENT)
    add_textbox(slide, customer,
                Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.5),
                font_size=22, bold=True, color=COLOR_DARK)
    for i, b in enumerate(bullets):
        add_textbox(slide, f"・{b}",
                    Inches(1.2), Inches(2.5) + Inches(0.55) * i, Inches(11), Inches(0.5),
                    font_size=20, color=COLOR_GRAY)
    # 下段：つながりたい方
    add_rect(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.45), RGBColor(0xFF, 0xF3, 0xE0))
    add_textbox(slide, connect_title,
                Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.45),
                font_size=18, bold=True, color=RGBColor(0xE6, 0x5C, 0x00))
    for i, c in enumerate(connects):
        add_textbox(slide, f"▶ {c}",
                    Inches(1.2), Inches(4.9) + Inches(0.55) * i, Inches(11), Inches(0.5),
                    font_size=20, color=COLOR_DARK)
    if note:
        slide.notes_slide.notes_text_frame.text = note


# ===== スライド生成 =====
prs = new_prs()

# --- SLIDE 01 タイトル ---
s = blank_slide(prs)
set_bg(s, COLOR_DARK)
add_textbox(s, '日本の“いいもの”を',
            Inches(1), Inches(1.8), Inches(11), Inches(1.0),
            font_size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(s, "ECショップで",
            Inches(1), Inches(2.8), Inches(11), Inches(1.0),
            font_size=48, bold=True, color=RGBColor(0x7B, 0xC8, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(s, "24時間、365日、全世界へ",
            Inches(1), Inches(3.8), Inches(11), Inches(1.0),
            font_size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(s, "鈴木",
            Inches(10.5), Inches(6.5), Inches(2.5), Inches(0.6),
            font_size=20, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)

# --- SLIDE 02 自己紹介① ---
s = blank_slide(prs)
slide_section(s, "自己紹介①　── まずは笑いを取る", [
    "鈴木（ひさこん）",
    "",
    "身長　185 cm",
    "体重　105 kg　→　目標 95 kg（ダイエット中）",
    "",
    "1ヶ月前から「お尻工場（新宿クロスフィット）」に通い始めました",
    "ケトジェニックダイエット実施中",
],
note="みなさん、こんにちは。鈴木です。\nまず最初に、自己紹介がてら正直に申し上げます。\n身長185cm、体重105kgです。（間）\n……目標は95kgです。ダイエット中です。\n1ヶ月前から「お尻工場」という、新宿にあるクロスフィットジムに通い始めまして、ケトジェニックダイエットも同時進行中です。\nこの大きな体で10分間お付き合いいただけると幸いです。よろしくお願いします！")

# --- SLIDE 03 自己紹介② ---
s = blank_slide(prs)
slide_section(s, "自己紹介②　── 経歴、なかなか面白いんです", [
    "🎓 高専・大学（情報工学 / AI工学）",
    "　　↓　勉強に飽きて……",
    "🪖 陸上自衛隊（伊丹・無線基地建設・小隊長　4年）",
    "　　↓　ものをつくりたい",
    "💻 WEB専門学校（東京・半年猛勉強）",
    "　　↓　リクルートに技術提供する会社のリーダーに気に入られ",
    "🚀 フリーランス → リクナビ / リクナビNEXT 開発グループ（5年）",
],
note="高専と大学で情報工学と人工知能工学を勉強していたんですが、ある日ふと思いまして。「……勉強、飽きたな」と。（間）\nそれで陸上自衛隊に入隊しました。（間）　理由はそれだけです。\n兵庫県の伊丹で、無線基地の建設と小隊長を4年やりました。\n\nで、今度は「ものをつくりたい」という気持ちが出てきて、東京のWEB専門学校に入学しました。\n半年、勉強しながら就職活動をしていたら、リクルートに技術提供している会社のリーダーに気に入っていただきまして。\nそのままフリーランスとしてリクルートに入り、リクナビとリクナビNEXTの開発グループで5年ほど働きました。\n\n……なかなか面白い経歴でしょう？（笑）")

# --- SLIDE 04 WANISE ---
s = blank_slide(prs)
slide_section(s, "直近の仕事 ── WANISE（ワンジー）", [
    '「1着で、誰でも美しく。日本の"着る哲学"を、世界へ。」',
    "",
    "・日本のユニフォームが持つ「誰が着ても美しく見える」哲学を日常服に応用",
    "・1着で複数の着こなしが可能",
    "・サイズ・体型を問わず美しく着られる設計",
    "",
    "▶ 私はこのブランドのECを担当しています",
],
note='私が今まさにECを担当しているブランドが、「ワンジー」です。\n\n日本にはユニフォームという独自の文化があります。「誰が着ても美しく見える」「機能と美を両立させる」——ワンジーは、そのユニフォームの哲学を日常のファッションにそのまま落とし込んだブランドです。\n\n1着で何通りもの着こなしができて、どんな体型の方にもフィットする。サイズを気にしなくていい服というのは、実は世界中の人が求めているものだと思います。\n\nこのブランドのECを担当し、日本の"着る哲学"を世界に届けるのが、今の私のミッションです。')

# --- SLIDE 05 私がやっていること ---
s = blank_slide(prs)
slide_section(s, "自己紹介③　── 私がやっていること", [
    "ECショップ構築 / 運用",
    "プロジェクトマネジメント",
    "プロダクト開発",
    "販売コミュニティ開発",
    "",
    "▶ 「ECだけじゃない。事業を動かす仕組みを、まるごとつくります。」",
],
note="ECだけをやっているように思われることが多いのですが、実はもう少し幅広くやっています。プロジェクトマネジメント、プロダクト開発、販売コミュニティの開発まで。ECを入口に、事業全体を動かす仕組みをつくるのが私のスタイルです。")

# --- SLIDE 06 ビジネス哲学 ---
s = blank_slide(prs)
slide_section(s, "ビジネスをどう考えているか", [
    '「"いいもの"が、届いていない。それを変えたい。」',
    "",
    "・日本には、世界に通用する商品・サービスを持つ人がたくさんいる",
    '・"売れない"のは商品のせいじゃない。届け方の問題',
    "・私がやりたいのは、その「仕組みのなさ」を解決すること",
],
note='本音をお話しすると——日本には、本当に素晴らしい商品やサービスを持っている方がたくさんいます。でも、その多くが「売る仕組み」を持っていないだけで、世の中に届いていない。私はそれが、もったいないと思っています。\n\n"売れない"のは、商品が悪いからじゃない。届け方の問題なんです。\n\n私がやりたいのは、そういう方々の力になること。良いものが、ちゃんと必要な人に届く。それを仕組みとしてつくるのが、私のビジネスに対する一番の動機です。')

# --- SLIDE 07 比較表 ---
s = blank_slide(prs)
slide_comparison(s,
    title="2店舗目、ネットショップにしませんか？",
    lead="2店舗目、またはアップセルを考えているなら——",
    rows=[
        ("保証金＋内装で数百万〜", "初期費用ほぼゼロ"),
        ("営業時間内だけ",         "24時間 365日"),
        ("出店エリアに限定",        "全世界に届く"),
    ],
    note="ここで、BNIメンバーの皆さんに一つ提案させてください。\n\nビジネスがうまくいき始めて、「そろそろ2店舗目を」「もっと単価を上げたい」と考えている方、いらっしゃいますか？\n\n実店舗を増やすとなると、保証金と内装だけで数百万円。それが当たり前になっていませんか。ネットショップなら、その初期費用がほぼゼロです。しかも24時間365日、エリアを問わず全世界に向けて開いている。アップセル商品をオンラインで展開すれば、既存のお客様にも新しい顧客にも同時にアプローチできます。\n\nメリットは他にもたくさんありますが、まずこの3つだけでも、選択肢として考えてみてほしいんです。次の一手を打つ前に、一度私に話を聞かせてください。")
add_textbox(s,
    "「他にもいろいろメリットはありますが、まずこの3つだけでも、検討する価値があると思います。」",
    Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
    font_size=18, italic=True, color=COLOR_GRAY)

# --- SLIDE 08 コンタクトサークル ---
s = blank_slide(prs)
slide_contact(s,
    title="コンタクトサークル",
    customer_title="【紹介してほしいお客様】",
    customer="▶ 2店舗目・アップセルを検討中の経営者",
    bullets=[
        "「売上は伸びているが、次のステップが見えない」",
        "「2店舗目を出したいが初期費用がネックで…」",
        "「既存客にもっと買ってもらう方法を探している」",
    ],
    connect_title="【今日つながりたい方】",
    connects=[
        "印刷・プリントのお仕事をされている方（宮原さん）",
        "グラフィックデザイナー",
    ],
    note="今日お願いしたいことが2つあります。\n\nまず、お客様のご紹介です。ビジネスが好調で「次の一手」を考えている経営者の方——2店舗目を検討している方、既存客へのアップセルを考えている方——そういう方がいたら、ぜひ紹介してください。\n\nそしてもう一つ。今日つながりたい方がいます。印刷・プリントのお仕事をされている宮原さん、そしてグラフィックデザイナーの方。ECサイトをつくる上で一緒に仕事をしたい方々です。心当たりがある方は、ぜひ声をかけてください。")

# --- SLIDE 09 デモ ---
s = blank_slide(prs)
slide_title_only(s, "デモ", "（別サイトを画面共有して見せる）")

# --- SLIDE 10 クロージング ---
s = blank_slide(prs)
set_bg(s, COLOR_DARK)
add_textbox(s, "詳しくは、1to1でお話しさせてください",
            Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2),
            font_size=34, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(s, "・2店舗目の前に、ECという選択肢を考えたい方\n・アップセルをオンラインで展開したい方\n・まずは話を聞いてみたい方、大歓迎です",
            Inches(2.5), Inches(3.0), Inches(8.5), Inches(2.0),
            font_size=22, color=RGBColor(0xCC, 0xDD, 0xFF))
add_textbox(s, "よろしくお願いいたします！",
            Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8),
            font_size=26, bold=True, color=RGBColor(0x7B, 0xC8, 0xFF), align=PP_ALIGN.CENTER)
s.notes_slide.notes_text_frame.text = "以上が私のプレゼンです。ありがとうございました。\n\n2店舗目を出す前にECという選択肢を考えてみたい方、アップセルをオンラインで展開したい方、ぜひ1to1でお話しさせてください。「話を聞くだけ」でも大歓迎です。よろしくお願いいたします！"

# ===== 保存 =====
out_path = "/Users/suzukiyoshihisa/Desktop/ 株式会社 HON/05_営業/BNI/メインプレゼン_10分.pptx"
prs.save(out_path)
print(f"✅ 保存完了: {out_path}")
